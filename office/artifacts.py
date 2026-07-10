"""Safe office artifact inspection, preservation checks, and optional rendering."""
from __future__ import annotations

import hashlib
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any


SUPPORTED = {".docx", ".xlsx", ".pptx", ".pdf", ".csv", ".txt", ".md"}


@dataclass
class ArtifactCheck:
    path: str
    kind: str
    ok: bool
    size: int
    sha256: str
    errors: list[str]
    warnings: list[str]
    pages_or_sheets: int = 0
    rendered_paths: list[str] | None = None

    def as_dict(self) -> dict[str, Any]:
        return asdict(self)


def _kind(path: str) -> str:
    return Path(path).suffix.lower().lstrip(".") or "unknown"


def _hash(path: str) -> str:
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def inspect_artifact(path: str, render: bool = False, render_dir: str | None = None) -> ArtifactCheck:
    path = os.path.realpath(os.path.expanduser(path))
    ext = Path(path).suffix.lower()
    errors: list[str] = []
    warnings: list[str] = []
    rendered: list[str] = []
    count = 0
    if not os.path.isfile(path):
        return ArtifactCheck(path, _kind(path), False, 0, "", ["file does not exist"], [], 0, [])
    if ext not in SUPPORTED:
        errors.append(f"unsupported artifact type: {ext or 'none'}")
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            count = len(doc.paragraphs) + len(doc.tables)
            if not count:
                warnings.append("document contains no paragraphs or tables")
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            count = len(prs.slides)
            if not count:
                errors.append("presentation contains no slides")
            for index, slide in enumerate(prs.slides, 1):
                if not any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes):
                    warnings.append(f"slide {index} has no readable text")
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=False, data_only=False)
            count = len(wb.worksheets)
            if not count:
                errors.append("workbook contains no worksheets")
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.startswith("=") and "#REF!" in cell.value:
                            errors.append(f"formula error in {ws.title}!{cell.coordinate}")
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            count = len(reader.pages)
            if not count:
                errors.append("PDF contains no pages")
        else:
            count = 1
            if os.path.getsize(path) == 0:
                warnings.append("artifact is empty")
    except Exception as exc:
        errors.append(f"parse failed: {exc}")
    if render and not errors and ext in {".docx", ".pptx", ".xlsx", ".pdf"}:
        rendered = render_artifact(path, render_dir)
        if not rendered:
            warnings.append("LibreOffice rendering unavailable; structural validation completed")
    return ArtifactCheck(path, _kind(path), not errors, os.path.getsize(path), _hash(path), errors, warnings, count, rendered)


def render_artifact(path: str, output_dir: str | None = None) -> list[str]:
    """Render through LibreOffice when installed; return generated PDF/image paths."""
    binary = shutil.which("libreoffice") or shutil.which("soffice")
    if not binary:
        return []
    output_dir = output_dir or tempfile.mkdtemp(prefix="captain-render-")
    os.makedirs(output_dir, exist_ok=True)
    result = subprocess.run([binary, "--headless", "--convert-to", "pdf", "--outdir", output_dir, path],
                            capture_output=True, text=True, timeout=90)
    if result.returncode != 0:
        return []
    pdf = os.path.join(output_dir, Path(path).stem + ".pdf")
    return [pdf] if os.path.isfile(pdf) else []


def preserve_unrelated_files(before: dict[str, str], after: dict[str, str]) -> list[str]:
    """Return paths whose content changed outside the declared artifact."""
    return sorted(path for path, digest in before.items() if after.get(path) != digest)
