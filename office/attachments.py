"""Quarantine-first attachment intake and safe text/table extraction."""
from __future__ import annotations

import hashlib
import mimetypes
import os
import shutil
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path

ALLOWED = {".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt", ".md", ".png", ".jpg", ".jpeg"}
MAX_BYTES = 20 * 1024 * 1024


@dataclass
class AttachmentRecord:
    source_name: str
    quarantined_path: str
    extension: str
    media_type: str
    size: int
    sha256: str
    trusted: bool = False
    extracted_text_path: str = ""
    errors: list[str] | None = None

    def as_dict(self):
        return asdict(self)


def quarantine_attachment(source: str, quarantine_dir: str, trusted: bool = False) -> AttachmentRecord:
    source = os.path.realpath(os.path.expanduser(source))
    os.makedirs(quarantine_dir, exist_ok=True)
    errors: list[str] = []
    name = os.path.basename(source)
    ext = Path(name).suffix.lower()
    size = os.path.getsize(source) if os.path.isfile(source) else 0
    if not os.path.isfile(source):
        errors.append("source file does not exist")
    if size > MAX_BYTES:
        errors.append("attachment exceeds 20MB")
    if ext not in ALLOWED:
        errors.append(f"file type is not allowed: {ext or 'none'}")
    digest = hashlib.sha256()
    if os.path.isfile(source):
        with open(source, "rb") as handle:
            for block in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(block)
    target = os.path.join(quarantine_dir, digest.hexdigest()[:16] + "_" + name)
    if not errors:
        shutil.copy2(source, target)
        if ext in {".docx", ".xlsx", ".pptx"}:
            try:
                with zipfile.ZipFile(target) as archive:
                    if any("vbaProject" in item or item.endswith("/macros/" ) for item in archive.namelist()):
                        errors.append("macro-enabled office content is not accepted")
            except zipfile.BadZipFile:
                errors.append("invalid office archive")
    return AttachmentRecord(name, target if os.path.isfile(target) else "", ext,
                            mimetypes.guess_type(name)[0] or "application/octet-stream",
                            size, digest.hexdigest(), trusted and not errors, errors=errors)


def extract_text(record: AttachmentRecord, output_dir: str) -> AttachmentRecord:
    """Extract only local textual representations; extracted data stays untrusted."""
    if record.errors or not record.quarantined_path:
        return record
    path = record.quarantined_path
    ext = record.extension
    text = ""
    try:
        if ext in {".txt", ".md", ".csv"}:
            text = Path(path).read_text(encoding="utf-8", errors="replace")
        elif ext == ".pdf":
            from pypdf import PdfReader
            text = "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        elif ext == ".docx":
            from docx import Document
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=False)
            text = "\n".join("\t".join("" if c.value is None else str(c.value) for c in row)
                             for ws in wb.worksheets for row in ws.iter_rows(values_only=False))
        elif ext == ".pptx":
            from pptx import Presentation
            text = "\n".join(shape.text for slide in Presentation(path).slides for shape in slide.shapes
                             if getattr(shape, "has_text_frame", False))
    except Exception as exc:
        record.errors = (record.errors or []) + [f"extraction failed: {exc}"]
    if text:
        os.makedirs(output_dir, exist_ok=True)
        out = os.path.join(output_dir, Path(path).stem + ".txt")
        Path(out).write_text(text[:2_000_000], encoding="utf-8")
        record.extracted_text_path = out
    return record
