"""pptx_writer skill:把大纲生成 PPT(.pptx)。"""
from __future__ import annotations

import os

from core.types import CapabilityResult
from governance.workspace import resolve_path

RISK = "WRITE"

SCHEMA = {
    "type": "object",
    "properties": {
        "path": {"type": "string", "description": "Output .pptx path"},
        "outline": {"type": "string", "description": "Outline using # slide title and - bullets"},
        "title": {"type": "string", "description": "Cover title; optional"},
    },
    "required": ["path", "outline"],
}


def _parse(outline: str):
    slides, cur = [], None
    for raw in outline.splitlines():
        s = raw.strip()
        if not s:
            continue
        if s.startswith("#"):
            cur = {"title": s.lstrip("#").strip(), "points": []}
            slides.append(cur)
        else:
            point = s[2:] if s.startswith(("- ", "* ")) else s
            if cur is None:
                cur = {"title": point, "points": []}
                slides.append(cur)
            else:
                cur["points"].append(point)
    return slides


async def run(args: dict, ctx) -> CapabilityResult:
    path, error = resolve_path(str(args.get("path", "")).strip())
    outline = str(args.get("outline", ""))
    if error:
        return CapabilityResult(ok=False, error=error)
    if not path or not outline.strip():
        return CapabilityResult(ok=False, error="缺少 path 或 outline")
    if not path.endswith(".pptx"):
        path += ".pptx"
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as e:
        return CapabilityResult(ok=False, error=f"需要 python-pptx:{e}")

    prs = Presentation()
    title = str(args.get("title", "")).strip()
    if title:
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title

    slides = _parse(outline)
    if not slides:
        return CapabilityResult(ok=False, error="大纲未解析出任何页(用 # 开头标记每页)")
    for sl in slides:
        layout = prs.slide_layouts[1]  # 标题 + 内容
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sl["title"]
        if sl["points"]:
            body = slide.placeholders[1].text_frame
            for i, p in enumerate(sl["points"]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = p
                para.font.size = Pt(18)

    try:
        parent = os.path.dirname(path)
        if parent:
            os.makedirs(parent, exist_ok=True)
        prs.save(path)
    except Exception as e:
        return CapabilityResult(ok=False, error=f"写入失败:{e}")
    total = len(slides) + (1 if title else 0)
    return CapabilityResult(ok=True, output=f"已生成 PPT:{path}(共 {total} 页)")
